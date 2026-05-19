from PyPDF2 import PdfReader
import io
from docx import Document
from pptx import Presentation
import fitz  
import io

def generate_pdf_thumbnail(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page = doc.load_page(0)  # first page

    pix = page.get_pixmap()
    img_bytes = pix.tobytes("png")

    return img_bytes
def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""

    for page in reader.pages:
        text += page.extract_text() or ""

    return text

def generate_id(text: str) -> str:
    return (
        text.lower()
        .replace("&", "and")
        .replace("/", "_")
        .replace(" ", "_")
    )

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    text.append(content)

    return "\n".join(text)

def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    text = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text.append(" | ".join(row_text))

    return "\n".join(text)

def generate_pdf_thumbnail(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page = doc.load_page(0)  # first page

    pix = page.get_pixmap()
    img_bytes = pix.tobytes("png")

    doc.close()

    return img_bytes